"""Build the APEX-ORACLE ISRO BAH 2026 idea-submission deck (Deep Space Midnight theme)."""
import random
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR

# ---- palette ----
NAVY  = RGBColor(0x0B, 0x14, 0x37)
NAVY2 = RGBColor(0x15, 0x20, 0x4F)
NAVY3 = RGBColor(0x1E, 0x2C, 0x66)
CYAN  = RGBColor(0x41, 0xE0, 0xD0)
ICE   = RGBColor(0xCA, 0xDC, 0xFC)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
MUTED = RGBColor(0x9A, 0xAE, 0xD6)
HEAD  = "Calibri"
BODY  = "Calibri"

prs = Presentation()
prs.slide_width = Emu(9144000)
prs.slide_height = Emu(5143500)
BLANK = prs.slide_layouts[6]
SW, SH = 10.0, 5.625

def slide():
    return prs.slides.add_slide(BLANK)

def bg(s, color=NAVY):
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    r.fill.solid(); r.fill.fore_color.rgb = color; r.line.fill.background()
    r.shadow.inherit = False
    return r

def stars(s, seed, n=26):
    rnd = random.Random(seed)
    for _ in range(n):
        x = rnd.uniform(0.1, 9.9); y = rnd.uniform(0.1, 5.5)
        d = rnd.choice([1.0, 1.4, 1.8, 2.4])
        col = CYAN if rnd.random() < 0.25 else RGBColor(0x6E, 0x86, 0xC0)
        e = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Pt(d), Pt(d))
        e.fill.solid(); e.fill.fore_color.rgb = col; e.line.fill.background(); e.shadow.inherit = False

def card(s, x, y, w, h, fill=NAVY2, line=None, lw=1.25, radius=True):
    shp = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE,
                             Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line; shp.line.width = Pt(lw)
    shp.shadow.inherit = False
    try:
        shp.adjustments[0] = 0.08
    except Exception:
        pass
    return shp

def text(s, x, y, w, h, content, size=14, color=WHITE, bold=False, align=PP_ALIGN.LEFT,
         anchor=MSO_ANCHOR.TOP, font=BODY, sp_after=4, line_spacing=1.05):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    lines = content.split("\n")
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.space_after = Pt(sp_after); p.space_before = Pt(0)
        if line_spacing: p.line_spacing = line_spacing
        r = p.add_run(); r.text = ln
        r.font.size = Pt(size); r.font.bold = bold; r.font.name = font
        r.font.color.rgb = color
    return tb

def runs_line(s, x, y, w, h, parts, size=12, anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.LEFT):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    p = tf.paragraphs[0]; p.alignment = align
    for txt, col, bold in parts:
        r = p.add_run(); r.text = txt; r.font.size = Pt(size); r.font.bold = bold
        r.font.name = BODY; r.font.color.rgb = col
    return tb

def numcircle(s, x, y, d, num):
    c = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(d), Inches(d))
    c.fill.solid(); c.fill.fore_color.rgb = CYAN; c.line.fill.background(); c.shadow.inherit = False
    tf = c.text_frame; tf.margin_left=0; tf.margin_right=0; tf.margin_top=0; tf.margin_bottom=0
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = str(num); r.font.size = Pt(13); r.font.bold = True
    r.font.color.rgb = NAVY; r.font.name = HEAD
    return c

def arrow(s, x, y, w=0.34, h=0.34, color=CYAN):
    a = s.shapes.add_shape(MSO_SHAPE.CHEVRON, Inches(x), Inches(y), Inches(w), Inches(h))
    a.fill.solid(); a.fill.fore_color.rgb = color; a.line.fill.background(); a.shadow.inherit = False
    return a

def kicker(s, txt, x=0.55, y=0.42):
    text(s, x, y, 9, 0.3, txt.upper(), size=12, color=CYAN, bold=True)

def title(s, txt, x=0.55, y=0.72, size=30, w=9):
    text(s, x, y, w, 0.7, txt, size=size, color=WHITE, bold=True)

# ===================== SLIDE 1 — TITLE =====================
s = slide(); bg(s); stars(s, 1, 34)
text(s, 0.6, 1.15, 9, 0.3, "BHARATIYA ANTARIKSH HACKATHON 2026", size=13, color=CYAN, bold=True)
runs_line(s, 0.6, 1.6, 9.2, 1.0, [("APEX", CYAN, True), ("-ORACLE", WHITE, True)], size=52)
text(s, 0.62, 2.75, 8.7, 0.8,
     "Physics-aware, explainable AI for detecting exoplanets in noisy TESS light curves",
     size=18, color=ICE, bold=False, line_spacing=1.1)
# bottom info card
card(s, 0.6, 4.05, 8.8, 1.15, fill=NAVY2)
runs_line(s, 0.85, 4.2, 8.4, 0.34, [("Problem Statement   ", CYAN, True),
          ("AI-enabled Detection of Exoplanets from Noisy Astronomical Light Curves", WHITE, False)], size=12)
runs_line(s, 0.85, 4.62, 8.4, 0.34, [("Team   ", CYAN, True), ("[YOUR TEAM NAME]      ", WHITE, False),
          ("Team Leader   ", CYAN, True), ("[YOUR NAME]   ", WHITE, False),
          ("(solo entry)", MUTED, False)], size=12)

# ===================== SLIDE 2 — TEAM =====================
s = slide(); bg(s); stars(s, 2)
kicker(s, "The Team"); title(s, "A solo entry, full-stack ownership")
card(s, 0.55, 1.55, 8.9, 1.35, fill=NAVY2, line=CYAN, lw=1.25)
ov = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.85), Inches(1.85), Inches(0.75), Inches(0.75))
ov.fill.solid(); ov.fill.fore_color.rgb = CYAN; ov.line.fill.background(); ov.shadow.inherit=False
tf = ov.text_frame; tf.margin_left=0; tf.margin_right=0; tf.margin_top=0; tf.margin_bottom=0
p = tf.paragraphs[0]; p.alignment=PP_ALIGN.CENTER; r=p.add_run(); r.text="YN"
r.font.size=Pt(20); r.font.bold=True; r.font.color.rgb=NAVY
text(s, 1.85, 1.9, 7.3, 0.45, "[YOUR NAME]", size=22, color=WHITE, bold=True)
text(s, 1.87, 2.42, 7.3, 0.4, "Solo participant  -  ML engineering + astro-informatics", size=14, color=ICE)
text(s, 0.55, 3.2, 8.9, 0.3, "One person, every hat in the pipeline:", size=13, color=MUTED)
hats = ["Data & detrending", "Transformer modeling", "Bayesian physics",
        "False-positive vetting", "Explainability (XAI)", "Live demo & report"]
hx, hy, hw, hh = 0.55, 3.6, 2.85, 0.6
for i, hlabel in enumerate(hats):
    cx = hx + (i % 3) * 2.97; cy = hy + (i // 3) * 0.78
    card(s, cx, cy, hw, hh, fill=NAVY3)
    text(s, cx+0.2, cy+0.14, hw-0.3, 0.32, hlabel, size=12.5, color=WHITE, bold=True, anchor=MSO_ANCHOR.MIDDLE)

# ===================== SLIDE 3 — OPPORTUNITY / USP =====================
s = slide(); bg(s); stars(s, 3)
kicker(s, "The Opportunity"); title(s, "Why APEX-ORACLE is different")
cols = [
    ("Different from existing ideas",
     "Most entries are black-box CNNs that only flag a brightness dip. APEX-ORACLE also classifies the signal, returns physical parameters with confidence, and shows visual reasons for every call."),
    ("How it solves the problem",
     "GP detrending preserves shallow transits, a multimodal model plus centroid vetting rejects blended false positives, and a differentiable transit model recovers parameters - even single transits."),
    ("Unique selling point",
     "Detection + classification + Bayesian parameters + explainability in one validated pipeline, proven by injection-recovery tests. Accurate and fully transparent - never a black box."),
]
cw = 2.85; cx0 = 0.55; gap = 0.13
for i, (h, b) in enumerate(cols):
    cx = cx0 + i * (cw + gap)
    card(s, cx, 1.6, cw, 3.45, fill=NAVY2, line=(CYAN if i == 2 else None), lw=1.25)
    numcircle(s, cx+0.25, 1.85, 0.42, i+1)
    text(s, cx+0.25, 2.45, cw-0.5, 0.8, h, size=15, color=CYAN, bold=True, line_spacing=1.05)
    text(s, cx+0.25, 3.35, cw-0.5, 1.55, b, size=11.5, color=ICE, line_spacing=1.12)

# ===================== SLIDE 4 — FEATURES =====================
s = slide(); bg(s); stars(s, 4)
kicker(s, "Capabilities"); title(s, "What the solution does")
feats = [
    ("Transit-safe detrending", "GP + Savitzky-Golay removes noise, keeps 100 ppm dips"),
    ("4-class signal typing", "Transit, eclipsing binary, blend, or starspot"),
    ("U-vs-V disentanglement", "Attention learns transit vs binary geometry"),
    ("Differentiable parameters", "Period, depth, duration with calibrated bounds"),
    ("False-positive vetting", "Centroid / difference imaging kills blends"),
    ("Monotransit recovery", "Finds single-transit, long-period planets"),
    ("Explainable visuals", "Phase-folded curves with attention overlays"),
    ("Honest validation", "Injection-recovery completeness & SNR"),
]
fw, fh = 4.35, 0.78; fx0, fy0 = 0.55, 1.55; gx, gy = 0.2, 0.12
for i, (h, b) in enumerate(feats):
    cx = fx0 + (i % 2) * (fw + gx); cy = fy0 + (i // 2) * (fh + gy)
    card(s, cx, cy, fw, fh, fill=NAVY2)
    dot = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(cx+0.22), Inches(cy+0.27), Inches(0.22), Inches(0.22))
    dot.fill.solid(); dot.fill.fore_color.rgb = CYAN; dot.line.fill.background(); dot.shadow.inherit=False
    text(s, cx+0.62, cy+0.1, fw-0.8, 0.32, h, size=13, color=WHITE, bold=True)
    text(s, cx+0.62, cy+0.42, fw-0.8, 0.3, b, size=10.5, color=MUTED)

# ===================== SLIDE 5 — PROCESS FLOW =====================
s = slide(); bg(s); stars(s, 5)
kicker(s, "Process Flow"); title(s, "Nine-stage detection pipeline")
stages = [
    ("Ingestion", "lightkurve, astroquery"),
    ("Detrending", "wotan, celerite2"),
    ("Candidate + fold", "transitleastsquares"),
    ("Classification", "Transformer + XGBoost"),
    ("Parameter estimation", "diff. batman + sbi"),
    ("Vetting", "centroid, triceratops"),
    ("Explainability", "attention, matplotlib"),
    ("Validation", "injection-recovery"),
    ("Live demo", "Streamlit dashboard"),
]
cw, ch = 2.8, 1.02; x0, y0 = 0.5, 1.45; csx, csy = 2.95, 1.22
for i, (h, b) in enumerate(stages):
    cx = x0 + (i % 3) * csx; cy = y0 + (i // 3) * csy
    hl = i in (3, 4)
    card(s, cx, cy, cw, ch, fill=(NAVY3 if hl else NAVY2), line=(CYAN if hl else None), lw=1.2)
    numcircle(s, cx+0.18, cy+0.16, 0.34, i+1)
    text(s, cx+0.62, cy+0.16, cw-0.75, 0.35, h, size=12.5, color=WHITE, bold=True)
    text(s, cx+0.62, cy+0.55, cw-0.75, 0.3, b, size=10, color=CYAN if hl else MUTED)

# ===================== SLIDE 6 — MOCKUP (Streamlit) =====================
s = slide(); bg(s); stars(s, 6)
kicker(s, "Mock-up (optional)"); title(s, "Live inspector dashboard")
card(s, 0.55, 1.5, 8.9, 3.6, fill=NAVY2)
# top bar
card(s, 0.75, 1.7, 8.5, 0.5, fill=NAVY3)
text(s, 0.95, 1.78, 6, 0.34, "APEX-ORACLE  -  Exoplanet Inspector", size=12, color=WHITE, bold=True, anchor=MSO_ANCHOR.MIDDLE)
card(s, 6.0, 1.78, 1.9, 0.34, fill=NAVY2)
text(s, 6.05, 1.8, 1.8, 0.3, "TIC 307210830", size=10.5, color=ICE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
ab = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.0), Inches(1.78), Inches(1.05), Inches(0.34))
ab.fill.solid(); ab.fill.fore_color.rgb = CYAN; ab.line.fill.background(); ab.shadow.inherit=False
tf=ab.text_frame; tf.margin_top=0; tf.margin_bottom=0; p=tf.paragraphs[0]; p.alignment=PP_ALIGN.CENTER
r=p.add_run(); r.text="Analyze"; r.font.size=Pt(10.5); r.font.bold=True; r.font.color.rgb=NAVY
# left: plot card with light curve
card(s, 0.75, 2.35, 5.1, 2.6, fill=NAVY3)
text(s, 0.95, 2.45, 4.7, 0.3, "Phase-folded light curve + attention", size=10.5, color=ICE, bold=True)
# draw transit curve with connectors
ax, ay, aw, ah = 0.95, 2.9, 4.7, 1.85
base = ay + ah*0.32; floor = ay + ah*0.74
N = 44; pts = []
for k in range(N):
    fx = k/(N-1)
    if 0.42 < fx < 0.46: yy = base + (floor-base)*((fx-0.42)/0.04)
    elif 0.46 <= fx <= 0.54: yy = floor
    elif 0.54 < fx < 0.58: yy = floor - (floor-base)*((fx-0.54)/0.04)
    else: yy = base
    pts.append((ax + fx*aw, yy))
for k in range(N-1):
    cn = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(pts[k][0]), Inches(pts[k][1]),
                                Inches(pts[k+1][0]), Inches(pts[k+1][1]))
    cn.line.color.rgb = CYAN; cn.line.width = Pt(2.0); cn.shadow.inherit=False
# right: results
card(s, 6.0, 2.35, 3.25, 0.62, fill=NAVY3)
runs_line(s, 6.2, 2.45, 3.0, 0.42, [("TRANSIT", CYAN, True), ("   confidence 96%", WHITE, False)], size=13)
res = [("Period", "3.41 d  +/- 0.02"), ("Depth", "910 ppm  +/- 40"),
       ("Duration", "2.7 h  +/- 0.1"), ("Vetting", "on-target, FPP 0.4%")]
for i,(k,v) in enumerate(res):
    cy = 3.08 + i*0.47
    runs_line(s, 6.05, cy, 3.2, 0.4, [(k+":  ", MUTED, False), (v, WHITE, True)], size=11)

# ===================== SLIDE 7 — ARCHITECTURE =====================
s = slide(); bg(s); stars(s, 7)
kicker(s, "Architecture"); title(s, "APEX-ORACLE hybrid pipeline")
boxes = [
    ("Data & labels", "TESS / MAST,\nGaia DR3, TOI", None),
    ("Detrend + BLS/TLS", "transit-safe GP,\ncandidate search", "A"),
    ("Transformer", "multimodal\nclassifier", "B"),
    ("Diff. batman + SBI", "Bayesian\nparameters", "C"),
    ("Vetting + output", "centroid, FPP,\nXAI overlays", None),
]
bw, bh = 1.66, 1.75; y = 1.65; x0 = 0.32; step = 1.92
for i, (h, b, tag) in enumerate(boxes):
    cx = x0 + i*step
    hl = tag in ("B", "C")
    card(s, cx, y, bw, bh, fill=(NAVY3 if hl else NAVY2), line=(CYAN if hl else None), lw=1.3)
    text(s, cx+0.13, y+0.18, bw-0.26, 0.7, h, size=12.5, color=WHITE, bold=True, line_spacing=1.0)
    text(s, cx+0.13, y+0.92, bw-0.26, 0.7, b, size=10, color=ICE, line_spacing=1.05)
    if tag:
        badge = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(cx+bw-0.42), Inches(y+bh-0.42), Inches(0.3), Inches(0.3))
        badge.fill.solid(); badge.fill.fore_color.rgb = CYAN; badge.line.fill.background(); badge.shadow.inherit=False
        tf=badge.text_frame; tf.margin_top=0; tf.margin_bottom=0; p=tf.paragraphs[0]; p.alignment=PP_ALIGN.CENTER
        rr=p.add_run(); rr.text=tag; rr.font.size=Pt(11); rr.font.bold=True; rr.font.color.rgb=NAVY
    if i < len(boxes)-1:
        arrow(s, cx+bw-0.02, y+bh/2-0.17, 0.3, 0.34)
# infra base
card(s, 0.32, 3.75, 9.36, 0.62, fill=NAVY2, line=CYAN, lw=1.0)
runs_line(s, 0.55, 3.83, 9.0, 0.46, [("Reproducibility & infra:   ", CYAN, True),
          ("uv  -  Hydra / OmegaConf  -  Weights & Biases  -  PyTorch Lightning  -  injection-recovery harness", WHITE, False)],
          size=11.5)
text(s, 0.32, 4.5, 9.3, 0.5, "A = fast front-end   .   B = detection spine   .   C = parameter & confidence head",
     size=11, color=MUTED)

# ===================== SLIDE 8 — TECHNOLOGIES =====================
s = slide(); bg(s); stars(s, 8)
kicker(s, "Technology Stack"); title(s, "Technologies used")
groups = [
    ("Data & astronomy", "lightkurve  -  astropy  -  astroquery\ntransitleastsquares  -  wotan"),
    ("Detrending & GP", "wotan  -  celerite2  -  gpytorch"),
    ("Deep learning", "PyTorch  -  Lightning  -  timm\nTransformers  -  XGBoost"),
    ("Physics & Bayesian", "batman  -  sbi  -  dynesty  -  emcee"),
    ("Vetting & XAI", "triceratops  -  attention maps\nmatplotlib  -  plotly"),
    ("Infra & delivery", "uv  -  Hydra  -  Weights & Biases\nStreamlit"),
]
gw, gh = 2.85, 1.45; gx0, gy0 = 0.55, 1.55; ggx, ggy = 0.15, 0.18
for i, (h, b) in enumerate(groups):
    cx = gx0 + (i % 3) * (gw + ggx); cy = gy0 + (i // 3) * (gh + ggy)
    card(s, cx, cy, gw, gh, fill=NAVY2)
    text(s, cx+0.22, cy+0.18, gw-0.4, 0.34, h, size=13, color=CYAN, bold=True)
    text(s, cx+0.22, cy+0.6, gw-0.4, 0.78, b, size=11, color=ICE, line_spacing=1.12)
text(s, 0.55, 5.05, 9, 0.3, "Language: Python 3.11   .   Compute: GPU for training, CPU-capable inference",
     size=11, color=MUTED)

# ===================== SLIDE 9 — COST =====================
s = slide(); bg(s); stars(s, 9)
kicker(s, "Estimated Cost (optional)"); title(s, "Near-zero, fully reproducible")
items = [
    ("Data", "TESS archive (MAST)", "Free"),
    ("Software", "100% open-source stack", "Free"),
    ("Compute", "Kaggle / Colab GPU, or short cloud-GPU rental", "Rs 0 - 8,000"),
]
for i, (k, d, v) in enumerate(items):
    cy = 1.7 + i*0.82
    card(s, 0.55, cy, 6.0, 0.7, fill=NAVY2)
    text(s, 0.8, cy+0.1, 1.7, 0.5, k, size=14, color=CYAN, bold=True, anchor=MSO_ANCHOR.MIDDLE)
    text(s, 2.4, cy+0.1, 4.0, 0.5, d, size=12, color=ICE, anchor=MSO_ANCHOR.MIDDLE)
    text(s, 5.4, cy+0.1, 1.05, 0.5, v, size=12.5, color=WHITE, bold=True, anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.RIGHT)
card(s, 6.85, 1.7, 2.6, 2.34, fill=NAVY3, line=CYAN, lw=1.3)
text(s, 6.95, 2.05, 2.4, 0.3, "TOTAL", size=12, color=MUTED, align=PP_ALIGN.CENTER, bold=True)
text(s, 6.95, 2.45, 2.4, 0.9, "~ Rs 0", size=46, color=CYAN, bold=True, align=PP_ALIGN.CENTER)
text(s, 6.95, 3.45, 2.4, 0.5, "to a few thousand", size=12, color=ICE, align=PP_ALIGN.CENTER)
text(s, 0.55, 4.75, 8.8, 0.5,
     "Low-risk and deployable: open data, open tools, runs on free-tier GPUs - the cost is engineering effort, not capital.",
     size=12, color=MUTED, line_spacing=1.1)

# ===================== SLIDE 10 — CLOSING =====================
s = slide(); bg(s); stars(s, 10, 40)
text(s, 0.6, 1.7, 9, 0.4, "APEX-ORACLE", size=20, color=CYAN, bold=True, align=PP_ALIGN.CENTER)
text(s, 0.6, 2.25, 8.8, 1.0, "Detect.  Explain.  Trust.", size=44, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
text(s, 0.6, 3.45, 8.8, 0.6,
     "Finding new worlds in the noise - and proving every discovery.",
     size=16, color=ICE, align=PP_ALIGN.CENTER)
text(s, 0.6, 4.7, 8.8, 0.4, "Bharatiya Antariksh Hackathon 2026   .   [YOUR NAME]   .   [YOUR TEAM NAME]",
     size=12, color=MUTED, align=PP_ALIGN.CENTER)

prs.save(r"E:\isro\deck\APEX-ORACLE_ISRO_BAH2026.pptx")
print("saved", len(prs.slides.__iter__.__self__._sldIdLst), "slides")
