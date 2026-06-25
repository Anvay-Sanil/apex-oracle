"""Fill the original ISRO BAH 2026 idea template with APEX-ORACLE content, keeping its format."""
import os

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN, MSO_ANCHOR

SRC = r"C:\Users\User\Downloads\[Pub] ISRO BAH 2026 _ Idea Submission Template.pptx"
OUT = r"E:\isro\deck\ISRO_BAH2026_Idea_Submission_FILLED.pptx"

NAVY = RGBColor(0x15, 0x22, 0x4F)
DARK = RGBColor(0x2B, 0x2B, 0x2B)
GRAY = RGBColor(0x55, 0x55, 0x55)
FONT = "Calibri"

p = Presentation(SRC)
slides = list(p.slides)

def box_with_text(slide, contains):
    for sh in slide.shapes:
        if sh.has_text_frame and contains.lower() in sh.text_frame.text.lower():
            return sh
    return None

def append_value(slide, label, value):
    sh = box_with_text(slide, label)
    para = sh.text_frame.paragraphs[0]
    r = para.add_run(); r.text = "   " + value
    r.font.name = FONT; r.font.size = Pt(15); r.font.bold = True; r.font.color.rgb = NAVY

def fill_content(slide, heading, bullets, note=None, hsize=20, bsize=15):
    sh = box_with_text(slide, heading_match(slide))
    # resize box to a comfortable content area below the header bar, above the footer strip
    sh.left = Inches(0.4); sh.top = Inches(0.95); sh.width = Inches(9.2); sh.height = Inches(4.35)
    tf = sh.text_frame; tf.clear(); tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE; tf.vertical_anchor = MSO_ANCHOR.TOP
    p0 = tf.paragraphs[0]; p0.space_after = Pt(10)
    r = p0.add_run(); r.text = heading; r.font.name = FONT; r.font.size = Pt(hsize)
    r.font.bold = True; r.font.color.rgb = NAVY
    for b in bullets:
        par = tf.add_paragraph(); par.space_after = Pt(7); par.line_spacing = 1.12
        r = par.add_run(); r.text = "•  " + b
        r.font.name = FONT; r.font.size = Pt(bsize); r.font.color.rgb = DARK
    if note:
        par = tf.add_paragraph(); par.space_before = Pt(6)
        r = par.add_run(); r.text = note
        r.font.name = FONT; r.font.size = Pt(11.5); r.font.italic = True; r.font.color.rgb = GRAY

# heading_match maps each content slide to a unique substring of its original guidance text
_HMAP = {}
def heading_match(slide):
    return _HMAP[id(slide)]

def set_match(slide, sub):
    _HMAP[id(slide)] = sub

# ---------- Slide 1: title labels ----------
append_value(slides[0], "team name", "[YOUR TEAM NAME]")
append_value(slides[0], "team leader", "[YOUR NAME]")
append_value(slides[0], "problem statement",
             "AI-enabled Detection of Exoplanets from Noisy Astronomical Light Curves")

# ---------- Slide 2: team members ----------
# The template already ships a baked-in grid (Team Leader + Team Member-1/2/3, each with
# Name/College), which matches ISRO's 3-4 member rule. Leave it untouched for the user to
# fill with real names once the team is formed - do not overlay our own list.

# ---------- Slides 3-9: content ----------
set_match(slides[2], "opportunity")
fill_content(slides[2], "Opportunity and USP", [
    "Most existing entries are black-box detectors that only flag a dip in brightness. Ours also says what the signal is, measures it, and shows its reasoning.",
    "It works end to end: noise-robust detrending keeps faint transits, a classifier plus a centroid check tells real planets from look-alike binaries, and a physics model fits the period, depth, and duration.",
    "USP: detection, classification, parameters with confidence, and clear visual explanations, all in one pipeline you can actually inspect.",
    "Already validated on real TESS data: an AI model trained on real light curves beats classical vetting 0.95 vs 0.75 accuracy in 5-fold cross-validation (65 targets), flagging 90% of real false positives while missing zero planets.",
], bsize=13)

set_match(slides[3], "list of features")
fill_content(slides[3], "Features of the solution", [
    "Transit-safe detrending (Gaussian Process and Savitzky-Golay) that removes noise but keeps the 100 ppm dips.",
    "Four-way classification: planet transit, eclipsing binary, blend, or starspot.",
    "Learns the U-shape versus V-shape difference using attention.",
    "Estimates period, depth, and duration with calibrated error bars.",
    "Rejects blended false positives using centroid and difference imaging.",
    "Recovers single-transit, long-period planets that folding methods miss.",
    "Phase-folded plots with attention overlays that explain each decision.",
    "Honest validation through injection-recovery tests and SNR.",
], bsize=14)

set_match(slides[4], "process flow")
fill_content(slides[4], "Process flow", [
    "Nine stages: ingestion, detrending, candidate search and phase fold, classification, parameter estimation, vetting, explainability, validation, and a live demo.",
    "Ingestion pulls TESS light curves and labels with lightkurve and astroquery.",
    "A Gaussian Process detrends, then transitleastsquares finds candidates and folds them.",
    "The folded signal is classified, fitted with a physics model, and checked for false positives.",
    "The output is a class, parameters, a confidence level, and explanatory plots in a dashboard.",
], note="A full process-flow diagram can be added on this slide.")

set_match(slides[5], "wireframes")
fill_content(slides[5], "Validation on real TESS data", [
    "AI trained on real TESS light curves beats classical vetting: 0.95 vs 0.75 accuracy (5-fold CV, 65 catalogue targets), and recovers real planets like WASP-18b at P=0.94 d.",
], bsize=13)
_fig = r"E:\isro\outputs\validation_panel.png"
if os.path.exists(_fig):
    slides[5].shapes.add_picture(_fig, Inches(0.5), Inches(2.15), width=Inches(9.0))

set_match(slides[6], "architecture diagram")
fill_content(slides[6], "Architecture", [
    "A hybrid of three parts: a fast candidate front-end, a Transformer classification spine, and a differentiable-physics parameter head.",
    "Data layer: TESS light curves and pixel data, Gaia neighbours, and TOI labels.",
    "Front-end: detrending and BLS/TLS produce clean transit candidates.",
    "Spine: a multimodal Transformer labels transit, eclipsing binary, blend, or starspot.",
    "Head: a differentiable batman model with simulation-based inference returns parameters and posteriors.",
    "Centroid and difference-imaging vetting removes false positives before the result is shown.",
    "Reproducibility layer: uv, Hydra, Weights and Biases, and PyTorch Lightning.",
], note="The architecture diagram can be added on this slide.", bsize=14)

set_match(slides[7], "technologies")
fill_content(slides[7], "Technologies used", [
    "Language and environment: Python 3.11 with uv.",
    "Data and astronomy: lightkurve, astropy, astroquery, transitleastsquares, wotan.",
    "Detrending and Gaussian Processes: wotan, celerite2, gpytorch.",
    "Deep learning: PyTorch, Lightning, timm, Transformers, XGBoost.",
    "Physics and Bayesian inference: batman, sbi, dynesty, emcee.",
    "Vetting and explainability: triceratops, attention maps, matplotlib, plotly.",
    "Tracking, config, and demo: Hydra, OmegaConf, Weights and Biases, Streamlit.",
], bsize=14)

set_match(slides[8], "estimated implementation cost")
fill_content(slides[8], "Estimated implementation cost (optional)", [
    "Data: the TESS archive from MAST is free.",
    "Software: the entire stack is open-source and free.",
    "Compute: free-tier GPUs (Kaggle or Colab), or a short cloud-GPU rental at roughly Rs 0 to 8,000.",
    "Total: close to zero. The real cost is engineering effort, not money.",
])

# ---------- Slide 10: closing ----------
# Native template "THANK YOU" slide sits on a dark image; leave it untouched (navy text
# would be invisible there and the closing is already complete).

p.save(OUT)
print("saved", OUT)
